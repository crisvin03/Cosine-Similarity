"""
Script to create PowerPoint presentation for Cosine Similarity Algorithm
CC 104 - Data Structures and Algorithms Project
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx library not installed!")
    print("Please install it by running:")
    print("  pip install python-pptx")
    exit(1)


def create_presentation():
    """Create the PowerPoint presentation for Cosine Similarity project."""
    
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions (standard 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(31, 78, 121)  # Dark blue
    text_color = RGBColor(51, 51, 51)     # Dark gray
    accent_color = RGBColor(0, 102, 204)  # Blue
    
    # Helper function to add title and content
    def add_slide(title_text, content_items, is_title_slide=False):
        """Add a slide with title and bullet points."""
        if is_title_slide:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        # Add title
        if is_title_slide:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = title_color
            title_para.alignment = PP_ALIGN.CENTER
        else:
            title = slide.shapes.title
            title.text = title_text
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = title_color
        
        # Add content
        if not is_title_slide and content_items:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(content_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Clean up markdown formatting
                clean_item = item.replace('**', '').replace('✅', '✓').replace('⚠️', '•')
                p.text = clean_item
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
                p.level = 0
                p.space_after = Pt(6)
        
        return slide
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Application of Cosine Similarity Algorithm\nin Information Retrieval and Text Analysis"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = title_color
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "CC 104 - Data Structures and Algorithms\nSorsogon State University - Bulan Campus\nDecember 2025"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = text_color
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Introduction
    add_slide(
        "Introduction",
        [
            "What is Cosine Similarity?",
            "Mathematical Formula:",
            "  cosine_similarity(A, B) = (A · B) / (||A|| × ||B||)",
            "Why is it important?",
            "Applications overview"
        ]
    )
    
    # Slide 3: Algorithm Overview
    add_slide(
        "Algorithm Overview",
        [
            "Core Components:",
            "  • Dot Product Calculation",
            "  • Vector Magnitude (Euclidean Norm)",
            "  • Cosine Similarity Formula",
            "",
            "Key Properties:",
            "  • Scale-invariant",
            "  • Range: -1 to 1",
            "  • Angle-based measurement"
        ]
    )
    
    # Slide 4: Mathematical Foundation
    add_slide(
        "Mathematical Foundation",
        [
            "Vector Space Model",
            "Term-Document Matrix",
            "TF-IDF Weighting",
            "Visual representation of vectors in space"
        ]
    )
    
    # Slide 5: Implementation Overview
    add_slide(
        "Implementation Overview",
        [
            "Programming Language: Python",
            "",
            "Key Features:",
            "  • Numerical vector similarity",
            "  • Text document similarity",
            "  • Recommendation system example",
            "",
            "Code Structure:",
            "  • CosineSimilarity class",
            "  • Modular design",
            "  • Well-documented"
        ]
    )
    
    # Slide 6: Demonstration 1
    add_slide(
        "Demonstration 1 - Numerical Vectors",
        [
            "Example 1: Identical vectors (similarity = 1.0)",
            "Example 2: Orthogonal vectors (similarity = 0.0)",
            "Example 3: Proportional vectors (similarity = 1.0)",
            "Example 4: Different vectors (similarity = 0.9746)",
            "",
            "Key Insight: Scale-invariance property"
        ]
    )
    
    # Slide 7: Demonstration 2
    add_slide(
        "Demonstration 2 - Text Document Similarity",
        [
            "Problem: Comparing text documents",
            "Solution: Convert text to vectors using term frequency",
            "Example: Comparing 4 documents about AI/ML and Python",
            "Results: Similarity matrix showing pairwise comparisons",
            "",
            "Application: Document clustering, plagiarism detection"
        ]
    )
    
    # Slide 8: Demonstration 3
    add_slide(
        "Demonstration 3 - Recommendation System",
        [
            "Problem: Recommend items to users",
            "Solution: Find similar users using cosine similarity",
            "Example: Movie recommendation based on user ratings",
            "Results: Identified most similar user (Eve, similarity = 0.9845)",
            "",
            "Application: E-commerce, content platforms"
        ]
    )
    
    # Slide 9: Applications
    add_slide(
        "Applications in Real World",
        [
            "1. Search Engines - Ranking search results",
            "2. Plagiarism Detection - Academic integrity",
            "3. Recommendation Systems - E-commerce, streaming",
            "4. Document Clustering - Organizing large collections",
            "5. Social Media - Content recommendation, trend analysis"
        ]
    )
    
    # Slide 10: Advantages
    add_slide(
        "Advantages",
        [
            "✓ Scale-invariant (handles different document lengths)",
            "✓ Computationally efficient",
            "✓ Works well with high-dimensional sparse vectors",
            "✓ Normalized output (-1 to 1)",
            "✓ Widely used and well-understood"
        ]
    )
    
    # Slide 11: Challenges
    add_slide(
        "Challenges and Limitations",
        [
            "• High-dimensional sparsity issues",
            "• Semantic limitations (bag-of-words)",
            "• Computational complexity at scale",
            "• Domain-specific adaptations needed",
            "",
            "Solutions: Word embeddings, dimensionality reduction, distributed computing"
        ]
    )
    
    # Slide 12: Literature Review
    add_slide(
        "Literature Review Findings",
        [
            "Domain: Information Retrieval and Text Analysis",
            "",
            "Key Findings:",
            "  • Cosine similarity is fundamental in IR systems",
            "  • Effective for document similarity and clustering",
            "  • Widely used in recommendation systems",
            "  • Integration with modern NLP techniques (embeddings)",
            "",
            "Research Gaps: Multilingual applications, real-time processing, deep learning integration"
        ]
    )
    
    # Slide 13: Code Quality
    add_slide(
        "Code Quality and Documentation",
        [
            "Features:",
            "  • Well-documented with clear comments",
            "  • Modular class-based design",
            "  • Error handling",
            "  • Type hints for clarity",
            "  • Multiple demonstration examples",
            "",
            "Best Practices: Follows Python conventions, readable code"
        ]
    )
    
    # Slide 14: Results
    add_slide(
        "Results and Analysis",
        [
            "Numerical Vectors: Successfully demonstrates basic cosine similarity",
            "Text Documents: Effectively compares documents of varying lengths",
            "Recommendation System: Identifies similar users accurately",
            "Performance: Efficient computation, handles edge cases"
        ]
    )
    
    # Slide 15: Future Work
    add_slide(
        "Future Work and Recommendations",
        [
            "1. Integration with semantic embeddings (Word2Vec, BERT)",
            "2. Scalability improvements for large datasets",
            "3. Multilingual and cross-lingual applications",
            "4. Hybrid approaches combining multiple similarity metrics",
            "5. Real-time processing optimizations"
        ]
    )
    
    # Slide 16: Conclusion
    add_slide(
        "Conclusion",
        [
            "Cosine similarity is a fundamental algorithm in information retrieval",
            "Effective for various applications: search, recommendations, text analysis",
            "Our implementation demonstrates core concepts with practical examples",
            "Well-documented code ready for extension and improvement",
            "",
            "Key Takeaway: Simple yet powerful algorithm with wide-ranging applications"
        ]
    )
    
    # Slide 17: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Thank you for your attention!"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = title_color
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions and Discussion"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = accent_color
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    filename = "Cosine_Similarity_Presentation.pptx"
    prs.save(filename)
    print(f"[OK] Presentation created successfully: {filename}")
    print(f"[OK] Total slides: {len(prs.slides)}")
    return filename


if __name__ == "__main__":
    try:
        print("Creating PowerPoint presentation...")
        filename = create_presentation()
        print(f"\n[SUCCESS] Presentation saved as: {filename}")
        print("\nYou can now:")
        print("  1. Open the file in Microsoft PowerPoint")
        print("  2. Customize colors, fonts, and add images")
        print("  3. Add your group members' names to the title slide")
        print("  4. Add diagrams or code snippets as needed")
    except Exception as e:
        print(f"\n[ERROR] Error creating presentation: {e}")
        import traceback
        traceback.print_exc()

